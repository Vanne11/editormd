#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import argparse, os, sys, shutil, tempfile, subprocess
from dataclasses import dataclass
from typing import List, Optional, Tuple

from markdown_it import MarkdownIt
from markdown_it.token import Token
from PIL import Image

# ===== DOCX / PPTX =====
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH

from pptx import Presentation
from pptx.util import Inches as PPTInches, Pt as PPTPt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

# ===== ReportLab (PDF) =====
from reportlab.lib.pagesizes import A4, LETTER
from reportlab.lib.units import mm
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image as RLImage,
    PageBreak, Preformatted, ListFlowable, ListItem
)
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# =========================
# Configuración de estilo
# =========================
@dataclass
class StyleConfig:
    font_family: str = "DejaVu Sans"   # nombre lógico
    font_size_pt: int = 12
    line_height: float = 1.4
    page: str = "A4"                   # "A4" | "Letter"
    ttf_path: Optional[str] = None     # ruta a .ttf opcional para PDF

# =========================
# Utilidades generales
# =========================
def which(cmd: str) -> Optional[str]:
    return shutil.which(cmd)

def ensure_dir(path: str):
    os.makedirs(path, exist_ok=True)

def sanitize_filename(name: str) -> str:
    return "".join(c if c.isalnum() or c in "._- " else "_" for c in name)

def img_path_resolve(src: str, base_dir: str) -> str:
    return src if os.path.isabs(src) else os.path.join(base_dir, src)

def find_images_in_inline(tokens: List[Token]) -> List[Tuple[str, str]]:
    out = []
    if not tokens:
        return out
    for t in tokens:
        if t.type == "image":
            alt = t.content or (t.attrs.get("alt", "") if t.attrs else "")
            src = ""
            if t.attrs:
                for k, v in t.attrs.items():
                    if k == "src":
                        src = v
            out.append((alt, src))
        if t.children:
            out.extend(find_images_in_inline(t.children))
    return out

# =========================
# Mermaid renderer (opcional)
# =========================
class MermaidRenderer:
    def __init__(self, workdir: str, prefer_svg: bool = True):
        self.workdir = workdir
        ensure_dir(workdir)
        self.mmdc = which("mmdc")
        self.prefer_svg = prefer_svg

    def available(self) -> bool:
        return self.mmdc is not None

    def render(self, code: str, name_hint: str = "diagram") -> Optional[str]:
        if not self.available():
            return None
        base = sanitize_filename(name_hint)
        ext = "png" if not self.prefer_svg else "svg"
        in_path = os.path.join(self.workdir, f"{base}.mmd")
        out_path = os.path.join(self.workdir, f"{base}.{ext}")
        with open(in_path, "w", encoding="utf-8") as f:
            f.write(code)
        cmd = [self.mmdc, "-i", in_path, "-o", out_path]
        try:
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            # Para DOCX/PPTX, PNG suele ser más compatible: si salió SVG, convertimos a PNG
            if out_path.endswith(".svg"):
                # Convertir SVG → PNG si Pillow no soporta SVG (asegurar compatibilidad)
                # Intento con rsvg-convert o cairosvg si existieran; si no, devolvemos SVG.
                if which("rsvg-convert"):
                    png_path = out_path.replace(".svg", ".png")
                    subprocess.run(["rsvg-convert", "-o", png_path, out_path], check=True)
                    return png_path
                elif which("cairosvg"):
                    png_path = out_path.replace(".svg", ".png")
                    subprocess.run(["cairosvg", out_path, "-o", png_path], check=True)
                    return png_path
            return out_path
        except subprocess.CalledProcessError:
            return None

# =========================
# Markdown → tokens (AST)
# =========================
def parse_markdown(md_text: str) -> List[Token]:
    md = MarkdownIt("commonmark", {"html": False, "linkify": True, "typographer": True})
    return md.parse(md_text)

# =========================
# ---------- PDF ----------
# =========================
def register_font_if_needed(style: StyleConfig):
    """
    Si recibimos una ruta a TTF, la registramos y usamos ese nombre
    como familia para los estilos de ReportLab.
    """
    if style.ttf_path and os.path.exists(style.ttf_path):
        font_name = os.path.splitext(os.path.basename(style.ttf_path))[0]
        try:
            pdfmetrics.registerFont(TTFont(font_name, style.ttf_path))
            return font_name
        except Exception:
            return style.font_family
    return style.font_family

def build_pdf_story(tokens: List[Token], style: StyleConfig, base_dir: str, mermaid: MermaidRenderer):
    story = []
    # Page size
    pagesize = A4 if style.page.lower() == "a4" else LETTER
    # Styles
    font_name = register_font_if_needed(style)
    ss = getSampleStyleSheet()
    body = ParagraphStyle(
        "Body",
        parent=ss["Normal"],
        fontName=font_name,
        fontSize=style.font_size_pt,
        leading=max(style.font_size_pt + 2, int(style.font_size_pt * style.line_height)),
        spaceAfter=6,
    )
    mono = ParagraphStyle(
        "Code",
        parent=ss["Code"],
        fontName=font_name,
        fontSize=max(9, int(style.font_size_pt * 0.9)),
        leading=max(10, int(style.font_size_pt * 1.2)),
        backColor="#f6f8fa",
        borderPadding=(6,6,6,6),
    )

    hstyles = {
        1: ParagraphStyle("H1", parent=body, fontSize=style.font_size_pt+8, leading=int((style.font_size_pt+8)*1.2), spaceBefore=12, spaceAfter=8),
        2: ParagraphStyle("H2", parent=body, fontSize=style.font_size_pt+4, leading=int((style.font_size_pt+4)*1.2), spaceBefore=10, spaceAfter=6),
        3: ParagraphStyle("H3", parent=body, fontSize=style.font_size_pt+2, leading=int((style.font_size_pt+2)*1.2), spaceBefore=8, spaceAfter=4),
        4: ParagraphStyle("H4", parent=body, fontSize=style.font_size_pt+1, spaceBefore=6, spaceAfter=4),
        5: ParagraphStyle("H5", parent=body, fontSize=style.font_size_pt, spaceBefore=6, spaceAfter=4),
        6: ParagraphStyle("H6", parent=body, fontSize=style.font_size_pt, spaceBefore=6, spaceAfter=4),
    }

    def add_image(path: str, max_width_mm: float = 160):
        if not os.path.exists(path):
            return
        try:
            with Image.open(path) as im:
                w, h = im.size
            # Escalar al ancho máximo preservando aspecto
            max_w = max_width_mm * mm
            dpi_scale = 72.0  # ReportLab usa puntos
            # Si la imagen está en px, asumimos 72 ppp (simplificado)
            ratio = min(1.0, max_w / w)
            iw = w * ratio
            ih = h * ratio
            img = RLImage(path, width=iw, height=ih)
            story.append(img)
            story.append(Spacer(1, 6))
        except Exception:
            # Si Pillow falla, intentamos insertar "tal cual"
            story.append(RLImage(path, width=160*mm))
            story.append(Spacer(1, 6))

    i = 0
    list_stack = []  # pila de listas (ListFlowable en construcción): [("ul"|"ol", ListFlowable, level)]
    def close_lists_until(level_target: int = 0):
        while len(list_stack) > level_target:
            _, lf, _ = list_stack.pop()
            story.append(lf)
            story.append(Spacer(1, 4))

    while i < len(tokens):
        t = tokens[i]

        # Encabezados
        if t.type == "heading_open":
            level = int(t.tag[1])
            text = ""
            if i+1 < len(tokens) and tokens[i+1].type == "inline":
                text = tokens[i+1].content
            close_lists_until(0)
            style_h = hstyles.get(level, body)
            story.append(Paragraph(text, style_h))
            i += 3
            continue

        # Párrafos (texto + imágenes en línea)
        if t.type == "paragraph_open":
            text = ""
            imgs = []
            if i+1 < len(tokens) and tokens[i+1].type == "inline":
                text = tokens[i+1].content
                imgs = find_images_in_inline(tokens[i+1].children or [])
            close_lists_until(len(list_stack))  # no cierras lista aquí; solo evitas confusiones
            if text.strip():
                story.append(Paragraph(text, body))
                story.append(Spacer(1, 2))
            for alt, src in imgs:
                ip = img_path_resolve(src, base_dir)
                add_image(ip)
            i += 3
            continue

        # Listas
        if t.type in ("bullet_list_open", "ordered_list_open"):
            kind = "ul" if t.type == "bullet_list_open" else "ol"
            lf = ListFlowable([], bulletType="bullet" if kind=="ul" else "1")
            list_stack.append((kind, lf, 0))
            i += 1
            continue

        if t.type == "list_item_open":
            # recolectar contenido simple (primer inline/paragraph) como ítem
            text = ""
            j = i + 1
            depth = 1
            while j < len(tokens) and depth > 0:
                if tokens[j].type == "list_item_open":
                    depth += 1
                elif tokens[j].type == "list_item_close":
                    depth -= 1
                j += 1
            # extraer primer inline entre i..j
            for k in range(i+1, j):
                if tokens[k].type == "inline":
                    text = tokens[k].content
                    break
            if list_stack:
                kind, lf, _ = list_stack[-1]
                li = ListItem(Paragraph(text, body))
                lf.flowables.append(li)
            i = j
            continue

        if t.type in ("bullet_list_close", "ordered_list_close"):
            # cerrar la lista actual
            close_lists_until(len(list_stack)-1)
            i += 1
            continue

        # Código (fence)
        if t.type == "fence":
            info = (t.info or "").strip().split()
            lang = info[0] if info else ""
            code = t.content
            if lang.lower() == "mermaid":
                out = mermaid.render(code, name_hint=f"mermaid_{i}")
                if out:
                    add_image(out)
                else:
                    story.append(Preformatted(code, mono))
            else:
                story.append(Preformatted(code, mono))
            story.append(Spacer(1, 4))
            i += 1
            continue

        # Imagen suelta
        if t.type == "image":
            src = ""
            if t.attrs:
                for k, v in t.attrs.items():
                    if k == "src":
                        src = v
            ip = img_path_resolve(src, base_dir)
            add_image(ip)
            i += 1
            continue

        # Salto de página manual si el MD trae '---' como hr (opcional)
        if t.type == "hr":
            story.append(PageBreak())
            i += 1
            continue

        i += 1

    # cerrar listas abiertas
    close_lists_until(0)
    return story, pagesize

def render_pdf_reportlab(tokens: List[Token], style: StyleConfig, base_dir: str, out_path: str, mermaid: MermaidRenderer):
    story, pagesize = build_pdf_story(tokens, style, base_dir, mermaid)
    doc = SimpleDocTemplate(
        out_path,
        pagesize=pagesize,
        leftMargin=20*mm, rightMargin=20*mm, topMargin=18*mm, bottomMargin=18*mm
    )
    doc.build(story)

# =========================
# ---------- DOCX ----------
# =========================
def render_docx(tokens: List[Token], style: StyleConfig, base_dir: str, out_path: str, mermaid: MermaidRenderer):
    doc = Document()
    st = doc.styles['Normal']
    st.font.name = style.font_family
    st.font.size = Pt(style.font_size_pt)

    i = 0
    while i < len(tokens):
        t = tokens[i]
        if t.type == "heading_open":
            level = int(t.tag[1])
            text = ""
            if i+1 < len(tokens) and tokens[i+1].type == "inline":
                text = tokens[i+1].content
            p = doc.add_paragraph()
            run = p.add_run(text)
            run.bold = True
            size_map = {1: style.font_size_pt+8, 2: style.font_size_pt+4, 3: style.font_size_pt+2}
            run.font.size = Pt(size_map.get(level, style.font_size_pt+1))
            i += 3
            continue

        if t.type == "paragraph_open":
            text = ""
            imgs = []
            if i+1 < len(tokens) and tokens[i+1].type == "inline":
                text = tokens[i+1].content
                imgs = find_images_in_inline(tokens[i+1].children or [])
            if text.strip():
                doc.add_paragraph(text)
            for alt, src in imgs:
                ip = img_path_resolve(src, base_dir)
                if os.path.exists(ip):
                    doc.add_picture(ip, width=Inches(6.0))
            i += 3
            continue

        if t.type == "fence":
            info = (t.info or "").strip().split()
            lang = info[0] if info else ""
            code = t.content
            if lang.lower() == "mermaid":
                out = mermaid.render(code, name_hint=f"mermaid_{i}")
                if out and os.path.exists(out):
                    doc.add_picture(out, width=Inches(6.0))
                else:
                    p = doc.add_paragraph()
                    r = p.add_run(code)
                    r.font.name = "Consolas"
                    r.font.size = Pt(max(9, int(style.font_size_pt*0.85)))
            else:
                p = doc.add_paragraph()
                r = p.add_run(code)
                r.font.name = "Consolas"
                r.font.size = Pt(max(9, int(style.font_size_pt*0.85)))
            i += 1
            continue

        if t.type == "image":
            src = ""
            if t.attrs:
                for k, v in t.attrs.items():
                    if k == "src":
                        src = v
            ip = img_path_resolve(src, base_dir)
            if os.path.exists(ip):
                doc.add_picture(ip, width=Inches(6.0))
            i += 1
            continue

        i += 1

    doc.save(out_path)

# =========================
# ---------- PPTX ----------
# =========================
def render_pptx(tokens: List[Token], style: StyleConfig, base_dir: str, out_path: str, mermaid: MermaidRenderer):
    prs = Presentation()
    title_layout = prs.slide_layouts[1]
    title_only = prs.slide_layouts[5]

    def new_slide(title_text: str):
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = title_text
        for p in slide.shapes.title.text_frame.paragraphs:
            for r in p.runs:
                r.font.name = style.font_family
                r.font.size = PPTPt(style.font_size_pt + 8)
        return slide

    def add_bullet(slide, text: str, level: int = 0):
        tf = slide.placeholders[1].text_frame
        if tf.text:
            p = tf.add_paragraph()
        else:
            p = tf.paragraphs[0]
        p.text = text
        p.level = min(level, 4)
        p.font.name = style.font_family
        p.font.size = PPTPt(style.font_size_pt)

    def add_image_slide(path: str, title: str = ""):
        slide = prs.slides.add_slide(title_only)
        if title:
            slide.shapes.title.text = title
        left = PPTInches(1)
        top = PPTInches(1.2)
        width = PPTInches(8)
        slide.shapes.add_picture(path, left, top, width=width)

    current = None
    i = 0
    while i < len(tokens):
        t = tokens[i]
        if t.type == "heading_open":
            level = int(t.tag[1])
            text = ""
            if i+1 < len(tokens) and tokens[i+1].type == "inline":
                text = tokens[i+1].content
            if level == 1 or current is None:
                current = new_slide(text or "Slide")
            else:
                if current is None:
                    current = new_slide("Contenido")
                add_bullet(current, text or "", level=level-1)
            i += 3
            continue

        if t.type == "paragraph_open":
            text = ""
            imgs = []
            if i+1 < len(tokens) and tokens[i+1].type == "inline":
                text = tokens[i+1].content
                imgs = find_images_in_inline(tokens[i+1].children or [])
            if imgs:
                for alt, src in imgs:
                    ip = img_path_resolve(src, base_dir)
                    if os.path.exists(ip):
                        add_image_slide(ip, title=alt or "")
            elif text.strip():
                if current is None:
                    current = new_slide("Contenido")
                add_bullet(current, text.strip(), level=0)
            i += 3
            continue

        if t.type == "fence":
            info = (t.info or "").strip().split()
            lang = info[0] if info else ""
            code = t.content
            if lang.lower() == "mermaid":
                out = mermaid.render(code, name_hint=f"mermaid_{i}")
                if out and os.path.exists(out):
                    add_image_slide(out, title="Diagrama")
                else:
                    if current is None:
                        current = new_slide("Contenido")
                    add_bullet(current, "[Mermaid no disponible]", level=0)
            else:
                if current is None:
                    current = new_slide("Contenido")
                add_bullet(current, "```" + lang + "```", level=0)
                for line in code.rstrip("\n").splitlines():
                    add_bullet(current, line, level=1)
            i += 1
            continue

        if t.type == "image":
            src = ""
            alt = t.content or ""
            if t.attrs:
                for k, v in t.attrs.items():
                    if k == "src":
                        src = v
            ip = img_path_resolve(src, base_dir)
            if os.path.exists(ip):
                add_image_slide(ip, title=alt or "")
            i += 1
            continue

        i += 1

    prs.save(out_path)

# =========================
# CLI
# =========================
def main():
    ap = argparse.ArgumentParser(description="Markdown → PDF (ReportLab), DOCX, PPTX. Tipografías, imágenes y Mermaid.")
    ap.add_argument("input", help="Ruta del .md de entrada")
    ap.add_argument("--to", choices=["pdf","docx","pptx"], required=True, help="Formato de salida")
    ap.add_argument("-o","--output", help="Ruta de salida (si se omite, se infiere)")
    ap.add_argument("--font-family", default="DejaVu Sans", help="Familia tipográfica (PDF/DOCX/PPTX)")
    ap.add_argument("--font-size", type=int, default=12, help="Tamaño de fuente (pt)")
    ap.add_argument("--line-height", type=float, default=1.4, help="Interlineado (PDF)")
    ap.add_argument("--page", choices=["A4","Letter"], default="A4", help="Tamaño de página PDF")
    ap.add_argument("--ttf-path", default=None, help="Ruta a un .ttf para PDF (se incrusta si es posible)")
    ap.add_argument("--mermaid-png", action="store_true", help="Forzar Mermaid en PNG (mejor compatibilidad)")
    args = ap.parse_args()

    in_path = args.input
    if not os.path.exists(in_path):
        print(f"ERROR: no existe {in_path}", file=sys.stderr); sys.exit(1)

    base_dir = os.path.abspath(os.path.dirname(in_path))
    md_text = open(in_path, "r", encoding="utf-8").read()
    tokens = parse_markdown(md_text)

    style = StyleConfig(
        font_family=args.font_family,
        font_size_pt=args.font_size,
        line_height=args.line_height,
        page=args.page,
        ttf_path=args.ttf_path
    )

    ext_map = {"pdf":".pdf","docx":".docx","pptx":".pptx"}
    out_path = args.output or (os.path.splitext(in_path)[0] + ext_map[args.to])

    with tempfile.TemporaryDirectory() as tmp:
        mermaid = MermaidRenderer(workdir=os.path.join(tmp, "mermaid"), prefer_svg=not args.mermaid_png)

        if args.to == "pdf":
            render_pdf_reportlab(tokens, style, base_dir, out_path, mermaid)
        elif args.to == "docx":
            render_docx(tokens, style, base_dir, out_path, mermaid)
        elif args.to == "pptx":
            render_pptx(tokens, style, base_dir, out_path, mermaid)

    print(f"✅ Listo: {out_path}")

if __name__ == "__main__":
    main()
